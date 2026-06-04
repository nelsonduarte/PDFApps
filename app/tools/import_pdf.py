"""PDFApps – TabImport: convert TXT, Images, Markdown, DOCX, PPTX, XLSX, HTML, EPUB to PDF."""

import os

from PySide6.QtCore import Qt
from PySide6.QtWidgets import (
    QGroupBox, QFormLayout, QComboBox, QLabel, QFileDialog,
    QMessageBox, QListWidget, QListWidgetItem,
    QAbstractItemView, QHBoxLayout, QPushButton,
)

from app.base import BasePage
from app.i18n import t
from app.utils import section, danger_btn, result_label_style
from app.constants import DESKTOP
from app.widgets import DropFileEdit


_IMG_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp", ".gif")


class TabImport(BasePage):
    def __init__(self, status_fn):
        super().__init__("fa5s.file-import", t("tool.import.name"),
                         t("tool.import.desc"),
                         t("tool.import.btn"), status_fn)
        f = self._form

        # ── Source type ──────────────────────────────────────────────
        grp = QGroupBox(t("tool.import.type_section"))
        gf = QFormLayout(grp)
        gf.setLabelAlignment(Qt.AlignmentFlag.AlignRight)
        self.cmb_type = QComboBox()
        self.cmb_type.addItems([
            t("tool.import.type_txt"),
            t("tool.import.type_images"),
            t("tool.import.type_md"),
            t("tool.import.type_docx"),
            t("tool.import.type_pptx"),
            t("tool.import.type_xlsx"),
            t("tool.import.type_html"),
            t("tool.import.type_epub"),
        ])
        self.cmb_type.currentIndexChanged.connect(self._on_type_changed)
        gf.addRow(t("tool.import.type_label"), self.cmb_type)
        f.addWidget(grp)

        # ── File list (works for all types) ──────────────────────────
        f.addWidget(section(t("tool.import.source_file")))

        self._file_list = QListWidget()
        self._file_list.setDragDropMode(QAbstractItemView.DragDropMode.InternalMove)
        self._file_list.setAlternatingRowColors(True)
        self._file_list.setMinimumHeight(140)
        f.addWidget(self._file_list)

        file_btns = QHBoxLayout()
        self._add_btn = QPushButton(t("tool.import.add_files"))
        self._add_btn.clicked.connect(self._pick_files)
        file_btns.addWidget(self._add_btn)
        self._remove_btn = danger_btn(t("btn.clear"))
        self._remove_btn.clicked.connect(self._file_list.clear)
        file_btns.addWidget(self._remove_btn)
        file_btns.addStretch()
        f.addLayout(file_btns)

        # ── Output ───────────────────────────────────────────────────
        f.addWidget(section(t("tool.import.output")))
        self.drop_out = DropFileEdit(save=True, default_name="output.pdf")
        f.addWidget(self.drop_out)

        self.lbl_result = QLabel("")
        self.lbl_result.setStyleSheet(result_label_style())
        f.addWidget(self.lbl_result)
        f.addStretch()

    def update_theme(self, dark: bool) -> None:
        super().update_theme(dark)
        try: self.lbl_result.setStyleSheet(result_label_style(dark))
        except RuntimeError: pass  # widget destroyed

    def _on_type_changed(self, index: int):
        self._file_list.clear()

    _FILTERS = {
        0: "tool.import.filter_txt",
        1: "file_filter.images",
        2: "tool.import.filter_md",
        3: "tool.import.filter_docx",
        4: "tool.import.filter_pptx",
        5: "tool.import.filter_xlsx",
        6: "tool.import.filter_html",
        7: "tool.import.filter_epub",
    }

    def _get_filter(self) -> str:
        key = self._FILTERS.get(self.cmb_type.currentIndex(), "tool.import.filter_txt")
        return t(key)

    def _pick_files(self):
        paths, _ = QFileDialog.getOpenFileNames(
            self, t("tool.import.add_files"), DESKTOP, self._get_filter())
        for p in paths:
            self._file_list.addItem(QListWidgetItem(p))
        if paths and not self.drop_out.path():
            base = os.path.splitext(paths[0])[0]
            self.drop_out.set_path(base + ".pdf")

    def _get_files(self) -> list:
        return [self._file_list.item(i).text() for i in range(self._file_list.count())]

    def _run(self):
        fmt = self.cmb_type.currentIndex()
        files = self._get_files()
        if not files:
            QMessageBox.warning(self, t("msg.warning"), t("tool.import.select_file"))
            return
        out = self._resolve_output_file(self.drop_out, files[0])
        if not out:
            return
        self.lbl_result.setText("")
        converters = {
            0: self._convert_txt,
            1: self._convert_images,
            2: self._convert_md,
            3: self._convert_docx,
            4: self._convert_pptx,
            5: self._convert_xlsx,
            6: self._convert_html,
            7: self._convert_epub,
        }
        converters[fmt](files, out)

    def _convert_txt(self, sources: list, out_path: str):
        n = len(sources)

        def do_work(worker):
            import fitz
            all_lines = []
            for i, src in enumerate(sources):
                if worker.is_cancelled():
                    return None
                with open(src, "r", encoding="utf-8") as f:
                    all_lines.extend(f.read().split("\n"))
                all_lines.append("")  # separator between files
                worker.progress.emit(i + 1, f"{i + 1}/{n}…")
            doc = fitz.open()
            page = None
            y = 50
            fontsize = 10
            line_height = fontsize * 1.4
            margin_x = 50
            max_y = 792
            max_width = 495  # 595 - 2*50
            for line in all_lines:
                if worker.is_cancelled():
                    doc.close()
                    return None
                if page is None or y + fontsize > max_y:
                    page = doc.new_page(width=595, height=842)  # A4
                    y = 50
                if not line.strip():
                    y += line_height
                    continue
                rect = fitz.Rect(margin_x, y, margin_x + max_width, max_y)
                used = page.insert_textbox(rect, line, fontsize=fontsize,
                                           fontname="helv")
                if used < 0:
                    page = doc.new_page(width=595, height=842)
                    y = 50
                    rect = fitz.Rect(margin_x, y, margin_x + max_width, max_y)
                    used = page.insert_textbox(rect, line, fontsize=fontsize,
                                               fontname="helv")
                est_lines = max(1, len(line) * fontsize * 0.5 / max_width + 1)
                y += line_height * est_lines
            try:
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    def _convert_images(self, sources: list, out_path: str):
        n = len(sources)

        def do_work(worker):
            import fitz
            doc = fitz.open()
            skipped = 0
            try:
                for i, img_path in enumerate(sources):
                    if worker.is_cancelled():
                        return None
                    img = fitz.open(img_path)
                    try:
                        if img.page_count == 0:
                            skipped += 1
                            continue
                        rect = img[0].rect
                        page = doc.new_page(width=rect.width, height=rect.height)
                        page.insert_image(page.rect, filename=img_path)
                    finally:
                        img.close()
                    worker.progress.emit(i + 1, f"{i + 1}/{n}…")
                doc.save(out_path)
            finally:
                doc.close()
            return skipped

        def on_done(skipped):
            if skipped:
                self._status(t("tool.import.skipped_images", n=skipped))
            self._done(out_path)

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=on_done)

    def _convert_md(self, sources: list, out_path: str):
        n = len(sources)

        def do_work(worker):
            import fitz
            all_md = []
            for i, src in enumerate(sources):
                if worker.is_cancelled():
                    return None
                with open(src, "r", encoding="utf-8") as f:
                    all_md.append(f.read())
                worker.progress.emit(i + 1, f"{i + 1}/{n}…")
            md_text = "\n\n---\n\n".join(all_md)
            lines = self._md_to_lines(md_text)
            doc = fitz.open()
            try:
                chunk = 55
                for i in range(0, max(len(lines), 1), chunk):
                    if worker.is_cancelled():
                        return None
                    page = doc.new_page(width=595, height=842)  # A4
                    y = 50
                    for text, size, bold in lines[i:i + chunk]:
                        if y > 780:
                            break
                        font = "helv" if not bold else "hebo"
                        try:
                            page.insert_text(fitz.Point(50, y), text,
                                             fontsize=size, fontname=font)
                        except Exception:
                            page.insert_text(fitz.Point(50, y), text,
                                             fontsize=size, fontname="helv")
                        y += size * 1.5
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    def _md_to_lines(self, md: str) -> list:
        """Convert markdown to list of (text, fontsize, bold) tuples."""
        result = []
        for line in md.split("\n"):
            stripped = line.strip()
            if stripped.startswith("# "):
                result.append((stripped[2:], 18, True))
            elif stripped.startswith("## "):
                result.append((stripped[3:], 15, True))
            elif stripped.startswith("### "):
                result.append((stripped[4:], 13, True))
            elif stripped.startswith("#### "):
                result.append((stripped[5:], 12, True))
            elif stripped.startswith("- ") or stripped.startswith("* "):
                result.append(("  \u2022  " + stripped[2:], 10, False))
            elif stripped.startswith("```"):
                continue  # skip code fences
            elif stripped == "---" or stripped == "***":
                result.append(("\u2500" * 60, 8, False))
            elif stripped == "":
                result.append(("", 10, False))
            else:
                # Remove inline formatting markers
                clean = stripped.replace("**", "").replace("__", "")
                clean = clean.replace("*", "").replace("_", "")
                clean = clean.replace("`", "")
                result.append((clean, 10, False))
        return result

    # ── DOCX → PDF ──────────────────────────────────────────────────────

    def _convert_docx(self, sources: list, out_path: str):
        try:
            from docx import Document  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_docx"))
            return
        n = len(sources)

        def do_work(worker):
            from docx import Document
            import fitz
            doc = fitz.open()
            try:
                for i, src in enumerate(sources):
                    if worker.is_cancelled():
                        return None
                    dx = Document(src)
                    lines = []
                    for para in dx.paragraphs:
                        text = para.text.strip()
                        style = (para.style.name or "").lower()
                        if "heading 1" in style:
                            lines.append((text, 18, True))
                        elif "heading 2" in style:
                            lines.append((text, 15, True))
                        elif "heading 3" in style:
                            lines.append((text, 13, True))
                        elif "heading" in style:
                            lines.append((text, 12, True))
                        elif text:
                            bold = any(r.bold for r in para.runs if r.bold)
                            lines.append((text, 10, bold))
                        else:
                            lines.append(("", 10, False))
                    for table in dx.tables:
                        lines.append(("", 10, False))
                        for row in table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            lines.append(("  |  ".join(cells), 9, False))
                        lines.append(("", 10, False))
                    self._render_lines_to_doc(doc, lines)
                    worker.progress.emit(i + 1, f"{i + 1}/{n}…")
                if worker.is_cancelled():
                    return None
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    # ── PPTX → PDF ──────────────────────────────────────────────────────

    def _convert_pptx(self, sources: list, out_path: str):
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_pptx"))
            return
        n = len(sources)

        def do_work(worker):
            from pptx import Presentation
            import fitz
            doc = fitz.open()
            try:
                for fi, src in enumerate(sources):
                    if worker.is_cancelled():
                        return None
                    prs = Presentation(src)
                    slide_w = prs.slide_width.pt if prs.slide_width else 960
                    slide_h = prs.slide_height.pt if prs.slide_height else 540
                    total_slides = len(prs.slides)
                    for i, slide in enumerate(prs.slides):
                        if worker.is_cancelled():
                            return None
                        page = doc.new_page(width=slide_w, height=slide_h)
                        y = 40
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if not text:
                                    continue
                                size = 10
                                bold = False
                                if para.runs:
                                    r = para.runs[0]
                                    if r.font.size:
                                        size = min(36, r.font.size.pt)
                                    bold = bool(r.font.bold)
                                font = "hebo" if bold else "helv"
                                if y + size > slide_h - 20:
                                    break
                                try:
                                    page.insert_text(fitz.Point(40, y), text,
                                                     fontsize=size, fontname=font)
                                except Exception:
                                    page.insert_text(fitz.Point(40, y), text,
                                                     fontsize=size, fontname="helv")
                                y += size * 1.4
                        # Slide-granular progress within the file; advance
                        # the dialog by (file index + slide fraction).
                        worker.progress.emit(
                            fi + 1,
                            f"{fi + 1}/{n}: {i + 1}/{total_slides}…")
                if worker.is_cancelled():
                    return None
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    # ── XLSX → PDF ──────────────────────────────────────────────────────

    def _convert_xlsx(self, sources: list, out_path: str):
        try:
            from openpyxl import load_workbook  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_xlsx"))
            return
        n = len(sources)

        def do_work(worker):
            from openpyxl import load_workbook
            import fitz
            doc = fitz.open()
            try:
                for i, src in enumerate(sources):
                    if worker.is_cancelled():
                        return None
                    wb = load_workbook(src, data_only=True)
                    for ws in wb.worksheets:
                        if worker.is_cancelled():
                            return None
                        lines = []
                        lines.append((f"— {ws.title} —", 12, True))
                        lines.append(("", 10, False))
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else "" for c in row]
                            lines.append(("  |  ".join(cells), 8, False))
                        lines.append(("", 10, False))
                        self._render_lines_to_doc(doc, lines)
                    worker.progress.emit(i + 1, f"{i + 1}/{n}…")
                if worker.is_cancelled():
                    return None
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    # ── HTML → PDF ──────────────────────────────────────────────────────

    def _convert_html(self, sources: list, out_path: str):
        try:
            from bs4 import BeautifulSoup  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.import.dep_bs4"))
            return
        n = len(sources)

        def do_work(worker):
            from bs4 import BeautifulSoup
            import fitz
            doc = fitz.open()
            try:
                for i, src in enumerate(sources):
                    if worker.is_cancelled():
                        return None
                    with open(src, "r", encoding="utf-8") as f:
                        soup = BeautifulSoup(f.read(), "html.parser")
                    lines = self._html_to_lines(soup)
                    self._render_lines_to_doc(doc, lines)
                    worker.progress.emit(i + 1, f"{i + 1}/{n}…")
                if worker.is_cancelled():
                    return None
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    def _html_to_lines(self, soup) -> list:
        lines = []
        tag_map = {"h1": (18, True), "h2": (15, True), "h3": (13, True),
                   "h4": (12, True), "h5": (11, True), "h6": (11, True)}
        for el in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6",
                                 "p", "li", "pre", "td", "th", "blockquote"]):
            text = el.get_text(strip=True)
            if not text:
                continue
            tag = el.name
            if tag in tag_map:
                size, bold = tag_map[tag]
                lines.append((text, size, bold))
            elif tag == "li":
                lines.append(("  \u2022  " + text, 10, False))
            elif tag == "pre":
                for ln in text.split("\n"):
                    lines.append((ln, 9, False))
            elif tag == "blockquote":
                lines.append(("  \u201c " + text + " \u201d", 10, False))
            elif tag in ("th",):
                lines.append((text, 10, True))
            else:
                lines.append((text, 10, False))
            if tag in tag_map:
                lines.append(("", 10, False))
        return lines

    # ── EPUB → PDF ──────────────────────────────────────────────────────

    def _convert_epub(self, sources: list, out_path: str):
        try:
            import ebooklib  # noqa: F401
            from ebooklib import epub  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_epub"))
            return
        try:
            from bs4 import BeautifulSoup  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.import.dep_bs4"))
            return
        n = len(sources)

        def do_work(worker):
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            import fitz
            doc = fitz.open()
            try:
                for i, src in enumerate(sources):
                    if worker.is_cancelled():
                        return None
                    book = epub.read_epub(src)
                    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                        if worker.is_cancelled():
                            return None
                        soup = BeautifulSoup(item.get_content(), "html.parser")
                        lines = self._html_to_lines(soup)
                        if lines:
                            self._render_lines_to_doc(doc, lines)
                    worker.progress.emit(i + 1, f"{i + 1}/{n}…")
                if worker.is_cancelled():
                    return None
                doc.save(out_path)
            finally:
                doc.close()
            return out_path

        self._run_background(do_work, total=max(n, 1),
                             label=t("tool.import.converting"),
                             on_done=lambda _r: self._done(out_path))

    # ── Shared line renderer ────────────────────────────────────────────

    def _render_lines_to_doc(self, doc, lines: list):
        """Render a list of (text, fontsize, bold) tuples to pages in an open fitz doc."""
        import fitz
        page = None
        y = 50
        for text, size, bold in lines:
            if page is None or y + size > 780:
                page = doc.new_page(width=595, height=842)
                y = 50
            if not text:
                y += size * 0.8
                continue
            font = "hebo" if bold else "helv"
            try:
                page.insert_text(fitz.Point(50, y), text,
                                 fontsize=size, fontname=font)
            except Exception:
                page.insert_text(fitz.Point(50, y), text,
                                 fontsize=size, fontname="helv")
            y += size * 1.5

    def _done(self, out_path: str):
        self.lbl_result.setText(f"  \u2192 {os.path.basename(out_path)}")
        self._status(t("tool.import.status.done", path=out_path))
        QMessageBox.information(self, t("msg.done"),
                                t("tool.import.done", path=out_path))
