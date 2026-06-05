"""PDFApps – TabConverter: convert PDF to images, DOCX, TXT, PPTX, XLSX, HTML, EPUB."""

import contextlib
import logging
import os
import re

_log = logging.getLogger(__name__)

_CTRL_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]')


def _clean(text: str) -> str:
    """Strip control characters that break XML-based formats."""
    return _CTRL_RE.sub('', text)

from PySide6.QtCore import Qt
from PySide6.QtWidgets import (
    QGroupBox, QFormLayout, QComboBox, QLabel, QFileDialog,
    QMessageBox, QProgressDialog,
)

from app.base import BasePage
from app.i18n import t
from app.utils import section, info_lbl, pick_folder, show_error, result_label_style
from app.constants import DESKTOP
from app.widgets import DropFileEdit


class TabConverter(BasePage):
    _DPI_VALUES = [72, 150, 300]

    def __init__(self, status_fn):
        super().__init__("fa5s.exchange-alt", t("tool.convert.name"),
                         t("tool.convert.desc"),
                         t("tool.convert.btn"), status_fn)
        f = self._form

        # -- Source file --
        sec_src = section(t("tool.convert.source"))
        f.addWidget(sec_src)
        self.drop_in = DropFileEdit()
        try: self.drop_in.btn.clicked.disconnect()
        except RuntimeError: pass
        self.drop_in.btn.clicked.connect(self._pick_input)
        self.drop_in.path_changed.connect(self._load_input)
        self.lbl_info = info_lbl()
        f.addWidget(self.drop_in); f.addWidget(self.lbl_info)

        # -- Output format --
        grp_fmt = QGroupBox(t("tool.convert.format_section"))
        gf = QFormLayout(grp_fmt)
        gf.setLabelAlignment(Qt.AlignmentFlag.AlignRight)
        self.cmb_format = QComboBox()
        self.cmb_format.addItems([
            t("tool.convert.png"), t("tool.convert.jpg"),
            t("tool.convert.docx"), t("tool.convert.txt"),
            t("tool.convert.pptx"), t("tool.convert.xlsx"),
            t("tool.convert.html"), t("tool.convert.epub"),
        ])
        self.cmb_format.currentIndexChanged.connect(self._on_format_changed)
        gf.addRow(t("tool.convert.format_label"), self.cmb_format)
        f.addWidget(grp_fmt)

        # -- Image options (visible for PNG/JPG) --
        self._grp_dpi = QGroupBox(t("tool.convert.img_options"))
        gd = QFormLayout(self._grp_dpi)
        gd.setLabelAlignment(Qt.AlignmentFlag.AlignRight)
        self.cmb_dpi = QComboBox()
        self.cmb_dpi.addItems([
            t("tool.convert.dpi_72"), t("tool.convert.dpi_150"), t("tool.convert.dpi_300"),
        ])
        self.cmb_dpi.setCurrentIndex(1)
        gd.addRow(t("tool.convert.dpi_label"), self.cmb_dpi)
        f.addWidget(self._grp_dpi)

        # -- Output folder (images) --
        sec_out_folder = section(t("tool.convert.output_folder"))
        f.addWidget(sec_out_folder)
        self._sec_out_folder = sec_out_folder
        self._drop_folder = DropFileEdit(placeholder=t("tool.convert.folder_hint"))
        try: self._drop_folder.btn.clicked.disconnect()
        except RuntimeError: pass
        self._drop_folder.btn.clicked.connect(self._pick_folder)
        f.addWidget(self._drop_folder)

        # -- Output file (DOCX / TXT) --
        self._section_file = section(t("tool.convert.output_file"))
        self._section_file.setVisible(False)
        f.addWidget(self._section_file)
        self._drop_file = DropFileEdit(save=True, default_name="converted.docx")
        self._drop_file.setVisible(False)
        f.addWidget(self._drop_file)

        self.lbl_result = QLabel("")
        self.lbl_result.setStyleSheet(result_label_style())
        f.addWidget(self.lbl_result)
        f.addStretch()
        self._compact_hidden = [sec_src, self.drop_in, self.lbl_info]
        # Hide all output sections — save dialog prompts automatically
        for w in (sec_out_folder, self._drop_folder, self._section_file, self._drop_file):
            w.setVisible(False)

    def update_theme(self, dark: bool) -> None:
        super().update_theme(dark)
        try: self.lbl_result.setStyleSheet(result_label_style(dark))
        except RuntimeError: pass  # widget destroyed

    # ── UI callbacks ──────────────────────────────────────────────────────

    _EXT_MAP = {2: ".docx", 3: ".txt", 4: ".pptx", 5: ".xlsx", 6: ".html", 7: ".epub"}

    def _on_format_changed(self, index: int):
        is_image = index <= 1
        self._grp_dpi.setVisible(is_image)
        if not self._compact_active:
            self._drop_folder.setVisible(is_image)
            self._sec_out_folder.setVisible(is_image)
            self._section_file.setVisible(not is_image)
            self._drop_file.setVisible(not is_image)
        if not is_image:
            ext = self._EXT_MAP.get(index, ".pdf")
            inp = self.drop_in.path()
            if inp:
                base = os.path.splitext(inp)[0]
                self._drop_file.set_path(base + ext)

    def _pick_input(self):
        p, _ = QFileDialog.getOpenFileName(self, t("btn.open_pdf"), DESKTOP, t("file_filter.pdf"))
        if p:
            self._load_input(p)

    def _load_input(self, p: str):
        self.drop_in.blockSignals(True)
        self.drop_in.set_path(p)
        self.drop_in.blockSignals(False)
        if not self._maybe_prompt_password(p):
            self.drop_in.blockSignals(True); self.drop_in.set_path("")
            self.drop_in.blockSignals(False); return
        size = os.path.getsize(p)
        try:
            r = self._open_reader(p)
            self.lbl_info.setText(t("tool.compress.pages_info", n=len(r.pages), size=f"{size/1024:.1f}"))
        except Exception as e:
            self.lbl_info.setText(t("tool.split.error_info", e=e))
        # auto-set output paths
        base = os.path.splitext(p)[0]
        if not self._drop_folder.path():
            self._drop_folder.blockSignals(True)
            self._drop_folder.set_path(os.path.dirname(p))
            self._drop_folder.blockSignals(False)
        fmt = self.cmb_format.currentIndex()
        if fmt >= 2:
            ext = self._EXT_MAP.get(fmt, ".pdf")
            self._drop_file.set_path(base + ext)

    def _pick_folder(self):
        d = pick_folder(self)
        if d:
            self._drop_folder.blockSignals(True)
            self._drop_folder.set_path(d)
            self._drop_folder.blockSignals(False)

    def auto_load(self, path: str):
        if path and not self.drop_in.path():
            self._load_input(path)

    # ── conversion logic ──────────────────────────────────────────────────

    def _run(self):
        pdf_path = self.drop_in.path()
        if not pdf_path or not os.path.isfile(pdf_path):
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return

        fmt = self.cmb_format.currentIndex()
        self.lbl_result.setText("")

        converters = {
            0: lambda p: self._convert_images(p, 0),
            1: lambda p: self._convert_images(p, 1),
            2: self._convert_docx,
            3: self._convert_txt,
            4: self._convert_pptx,
            5: self._convert_xlsx,
            6: self._convert_html,
            7: self._convert_epub,
        }
        converters[fmt](pdf_path)

    def _make_progress(self, total: int, label: str) -> QProgressDialog:
        progress = QProgressDialog(label, t("progress.cancel"), 0, total, self)
        progress.setWindowTitle(t("progress.compress.title"))
        progress.setWindowModality(Qt.WindowModality.WindowModal)
        progress.setMinimumDuration(0)
        progress.setValue(0)
        return progress

    def _convert_images(self, pdf_path: str, fmt: int):
        out_dir = self._resolve_output_dir(self._drop_folder, pdf_path)
        if not out_dir:
            return
        os.makedirs(out_dir, exist_ok=True)
        ext = "png" if fmt == 0 else "jpg"
        dpi = self._DPI_VALUES[self.cmb_dpi.currentIndex()]
        self._status(t("tool.convert.status.starting",
                       format=ext.upper(), dpi=dpi))

        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            show_error(self, e); return

        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                matrix = fitz.Matrix(dpi / 72, dpi / 72)
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                    pix = page.get_pixmap(matrix=matrix)
                    if pix.alpha:
                        pix = fitz.Pixmap(pix, 0)
                    if pix.n == 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    out_file = os.path.join(out_dir, f"page_{i + 1:03d}.{ext}")
                    if ext == "png":
                        pix.save(out_file)
                    else:
                        try:
                            from PIL import Image
                            mode = "L" if pix.n == 1 else "RGB"
                            img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
                            img.save(out_file, "JPEG", quality=95)
                        except ImportError:
                            pix.save(out_file)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  {result} → {out_dir}")
            self._status(t("tool.convert.status.images_done", n=result))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_images",
                                      n=result, folder=out_dir))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    def _convert_docx(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.docx")
        if not out_path:
            return
        # Pre-flight on main thread: dep checks + page count + capture pwd.
        try:
            import fitz  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.ocr.dep_pymupdf"))
            return
        try:
            from docx import Document  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_docx"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            show_error(self, e)
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            import io, re as _re
            from app.tools._pdf_extract import (
                extract_page_assets,
                detect_repeated_regions,
                detect_card_regions,
                detect_table_regions,
            )
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                # Pass 1: extract assets for every page (text blocks, images,
                # widgets, annotations) using the shared helper.
                pages_assets = []
                for i in range(doc.page_count):
                    if worker.is_cancelled():
                        return None
                    pages_assets.append(extract_page_assets(doc, i))
                    worker.progress.emit(i, f"{i + 1}/{total}…")

                # Pass 2: detect headers/footers that repeat across pages so
                # they are not duplicated in the body flow. The helper uses
                # each page's own height, so mixed portrait/landscape PDFs
                # are handled correctly; the second arg is just a fallback.
                skip = detect_repeated_regions(pages_assets)

                # Pass 3: write the DOCX out of the cached assets.
                docx_doc = Document()
                for pa in pages_assets:
                    if worker.is_cancelled():
                        return None
                    # 3a. Inline images (best-effort).
                    for img in pa.images:
                        if not img.bytes:
                            continue
                        with contextlib.suppress(Exception):
                            para = docx_doc.add_paragraph()
                            run = para.add_run()
                            run.add_picture(io.BytesIO(img.bytes), width=Inches(5.0))

                    # 3a-bis. Detect colored callout / card regions and
                    # rasterize them as inline PNG images at 200 DPI. Text
                    # blocks contained within a card are consumed (suppressed
                    # from the flow below) to avoid duplication. Tables —
                    # grids of filled cells — are excluded by the card
                    # detector and handled by ``detect_table_regions`` below
                    # so they survive as real ``<w:tbl>`` elements.
                    try:
                        card_regions = detect_card_regions(pa)
                    except Exception as exc:
                        _log.warning(
                            "detect_card_regions failed on page %d: %s",
                            pa.page_index, exc,
                        )
                        card_regions = []
                    # 3a-ter. Detect tabular grids — Phase E3. Cards and
                    # tables should be mutually exclusive (the card
                    # detector already drops grid-like regions) but a
                    # standalone "callout" can occasionally bbox-overlap
                    # a multi-cell layout that the table detector also
                    # picks up. We give tables priority and drop any
                    # card whose bbox is largely covered by a detected
                    # table region — that prevents the same text from
                    # appearing both inside a rasterized card image and
                    # again inside the table cell.
                    try:
                        table_regions = detect_table_regions(pa)
                    except Exception as exc:
                        _log.warning(
                            "detect_table_regions failed on page %d: %s",
                            pa.page_index, exc,
                        )
                        table_regions = []

                    def _bbox_overlap_ratio_local(
                        a: tuple[float, float, float, float],
                        b: tuple[float, float, float, float],
                    ) -> float:
                        ax0, ay0, ax1, ay1 = a
                        bx0, by0, bx1, by1 = b
                        ix0 = max(ax0, bx0); iy0 = max(ay0, by0)
                        ix1 = min(ax1, bx1); iy1 = min(ay1, by1)
                        if ix1 <= ix0 or iy1 <= iy0:
                            return 0.0
                        inter = (ix1 - ix0) * (iy1 - iy0)
                        a_area = max(1e-6, (ax1 - ax0) * (ay1 - ay0))
                        return inter / a_area

                    if table_regions:
                        # Drop cards whose bbox is mostly inside a table
                        # (>= 50 % of the card's own area). Mutual
                        # exclusion: text shows up only as table cells.
                        card_regions = [
                            cr for cr in card_regions
                            if not any(
                                _bbox_overlap_ratio_local(cr.bbox, tr.bbox) >= 0.5
                                for tr in table_regions
                            )
                        ]

                    consumed_text_indices: set[int] = set()
                    consumed_widget_indices: set[int] = set()
                    consumed_annotation_indices: set[int] = set()
                    for cr in card_regions:
                        for ti in cr.text_block_indices:
                            consumed_text_indices.add(ti)
                        for wi in cr.widget_indices:
                            consumed_widget_indices.add(wi)
                        for ai in cr.annotation_indices:
                            consumed_annotation_indices.add(ai)
                    for tr in table_regions:
                        for ti in tr.text_block_indices:
                            consumed_text_indices.add(ti)
                        # Phase E3 fix: widgets / annotations whose bbox
                        # is inside the table region are already rendered
                        # by the cell text, so suppress them from the
                        # trailing widget / annotation passes.
                        for wi in tr.widget_indices:
                            consumed_widget_indices.add(wi)
                        for ai in tr.annotation_indices:
                            consumed_annotation_indices.add(ai)

                    def _emit_card(cr) -> None:
                        try:
                            page = doc[pa.page_index]
                            # Defensive: page.get_drawings() returns coords in
                            # the post-rotation space and get_pixmap(clip=...)
                            # expects the same, so the existing flow already
                            # works for rotated pages. If a regression ever
                            # surfaces we may need to apply
                            # ``page.derotation_matrix`` to the clip rect.
                            # TODO(E3): explicit derotation if a real-world
                            # rotated PDF reproduces a misalignment.
                            if getattr(page, "rotation", 0):
                                _log.debug(
                                    "page %d rotated %d°, card clip may "
                                    "need derotation",
                                    pa.page_index, page.rotation,
                                )
                            clip = fitz.Rect(*cr.bbox)
                            pix = page.get_pixmap(
                                clip=clip,
                                dpi=200,
                                colorspace=fitz.csRGB,
                            )
                            data = pix.tobytes("png")
                            pix = None  # release native buffer
                            if not data:
                                return
                            page_w = pa.width or 1.0
                            card_w = max(0.0, cr.bbox[2] - cr.bbox[0])
                            # Word usable width ~ 6 inches (Letter, 1" margins).
                            width_in = (card_w / page_w) * 6.0
                            # Floor very low so narrow sidebar callouts are
                            # not blown up 2x; cap below the Letter usable
                            # width so a near-full-page card doesn't push
                            # past the right margin.
                            width_in = max(width_in, 1.0)
                            width_in = min(width_in, 6.5)
                            para = docx_doc.add_paragraph()
                            run = para.add_run()
                            run.add_picture(
                                io.BytesIO(data), width=Inches(width_in)
                            )
                        except Exception as exc:
                            _log.warning(
                                "card emit failed on page %d: %s",
                                pa.page_index, exc,
                            )
                            return

                    def _emit_table(tr) -> None:
                        try:
                            docx_table = docx_doc.add_table(
                                rows=tr.rows, cols=tr.cols
                            )
                            # Default Word style — produces visible borders.
                            # Borderless detection is recorded in
                            # ``tr.has_borders`` for future styling work.
                            with contextlib.suppress(Exception):
                                docx_table.style = "Table Grid"
                            for cell in tr.cells:
                                if cell.row >= tr.rows or cell.col >= tr.cols:
                                    continue
                                texts: list[str] = []
                                for ti in cell.text_block_indices:
                                    if ti < 0 or ti >= len(pa.text_blocks):
                                        continue
                                    block = pa.text_blocks[ti]
                                    for line in block.lines:
                                        line_txt = "".join(
                                            _clean(s.text) for s in line.spans
                                        ).strip()
                                        if line_txt:
                                            texts.append(line_txt)
                                try:
                                    docx_cell = docx_table.rows[cell.row].cells[
                                        cell.col
                                    ]
                                    # ``cell.text = "\n".join(...)`` injects a
                                    # literal '\n' character into a single
                                    # paragraph — Word does not interpret it
                                    # as a line break. Instead, seed the cell
                                    # with the first line and append the rest
                                    # as additional paragraphs so the visual
                                    # multi-line layout survives.
                                    docx_cell.text = texts[0] if texts else ""
                                    for extra in texts[1:]:
                                        docx_cell.add_paragraph(extra)
                                except Exception:
                                    continue
                        except Exception as exc:
                            _log.warning(
                                "table emit failed on page %d: %s",
                                pa.page_index, exc,
                            )

                    # Merge cards + tables into a single Y-ordered queue so
                    # the flush logic below preserves visual order when
                    # both kinds appear on the same page.
                    pending: list[tuple[float, str, object]] = []
                    for cr in card_regions:
                        pending.append((cr.bbox[1], "card", cr))
                    for tr in table_regions:
                        pending.append((tr.bbox[1], "table", tr))
                    pending.sort(key=lambda x: x[0])
                    pending_idx = 0

                    def _flush_until(y_limit: float) -> None:
                        nonlocal pending_idx
                        while pending_idx < len(pending):
                            top, kind, obj = pending[pending_idx]
                            if top > y_limit:
                                return
                            if kind == "card":
                                _emit_card(obj)
                            else:
                                _emit_table(obj)
                            pending_idx += 1

                    # 3b. Text blocks (skip repeated header/footer).
                    for bi, block in enumerate(pa.text_blocks):
                        if (pa.page_index, bi) in skip:
                            continue
                        # Emit any pending cards / tables whose top edge is
                        # above the current block — keeps visual order.
                        _flush_until(block.bbox[1])
                        if bi in consumed_text_indices:
                            continue
                        lines = block.lines
                        if not lines:
                            continue
                        all_spans = [s for line in lines for s in line.spans]
                        if not all_spans:
                            continue
                        block_text = " ".join(
                            _clean(s.text) for s in all_spans
                        ).strip()
                        if not block_text:
                            continue
                        # Skip standalone page numbers
                        if _re.match(r"^\s*(?:page\s+)?\d{1,4}(?:\s+of\s+\d{1,4})?\s*$",
                                     block_text, _re.IGNORECASE):
                            continue
                        # Skip TOC dot-leader lines
                        if _re.search(r'\.[\s.]*\.[\s.]*\.[\s.]*\.', block_text):
                            continue
                        # Detect heading level by font size
                        max_size = max((s.size or 12) for s in all_spans)
                        any_bold = any((s.flags & 16) for s in all_spans)
                        if max_size >= 20:
                            para = docx_doc.add_heading(level=1)
                        elif max_size >= 16:
                            para = docx_doc.add_heading(level=2)
                        elif max_size >= 13 and any_bold:
                            para = docx_doc.add_heading(level=3)
                        elif any_bold and max_size >= 11:
                            para = docx_doc.add_heading(level=4)
                        else:
                            para = docx_doc.add_paragraph()
                        for li, line in enumerate(lines):
                            for span in line.spans:
                                text = _clean(span.text)
                                if not text:
                                    continue
                                run = para.add_run(text)
                                run.font.size = Pt(span.size or 12)
                                sf = span.flags
                                run.font.bold = bool(sf & 16)
                                run.font.italic = bool(sf & 2)
                                color = span.color
                                if color and color != 0:
                                    r_val = (color >> 16) & 0xFF
                                    g_val = (color >> 8) & 0xFF
                                    b_val = color & 0xFF
                                    run.font.color.rgb = RGBColor(r_val, g_val, b_val)
                            if li < len(lines) - 1:
                                para.add_run(" ")

                    # Flush any remaining cards / tables that sat below the
                    # last text block (or pages whose only content is a
                    # card / table).
                    _flush_until(float("inf"))

                    # 3c. Form widgets — emit captured values so they are not lost.
                    for wi, w in enumerate(pa.widgets):
                        if wi in consumed_widget_indices:
                            # Already visible inside a rasterized card.
                            continue
                        value = _clean(w.field_value).strip()
                        if not value:
                            continue
                        name = _clean(w.field_name).strip() or w.field_type or "field"
                        docx_doc.add_paragraph(f"[Form: {name}] {value}")

                    # 3d. Annotations — sticky notes, FreeText, etc.
                    for ai, a in enumerate(pa.annotations):
                        if ai in consumed_annotation_indices:
                            # Already visible inside a rasterized card.
                            continue
                        content = _clean(a.content).strip()
                        if not content:
                            continue
                        docx_doc.add_paragraph(f"[Note: {a.type}] {content}")
                if worker.is_cancelled():
                    return None
                docx_doc.save(out_path)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(t("tool.convert.status.docx_done", path=out_path))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_docx", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    def _convert_txt(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.txt")
        if not out_path:
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            show_error(self, e)
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                with open(out_path, 'w', encoding='utf-8') as f:
                    for i, page in enumerate(doc):
                        if worker.is_cancelled():
                            return None
                        if i > 0:
                            f.write(t("tool.convert.txt.page_separator",
                                      n=i + 1))
                        f.write(page.get_text())
                        worker.progress.emit(i, f"{i + 1}/{total}…")
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(t("tool.convert.status.txt_done", path=out_path))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_txt", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → PPTX ──────────────────────────────────────────────────────

    def _convert_pptx(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.pptx")
        if not out_path:
            return
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_pptx"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
                first = _probe[0].rect if total else None
        except Exception as e:
            show_error(self, e)
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        slide_w_pt = first.width
        slide_h_pt = first.height
        pwd = self._pdf_password

        def do_work(worker):
            import fitz, io
            from pptx import Presentation
            from pptx.util import Emu, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.oxml.ns import qn
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)

            def _rgb(c):
                """fitz colors are 0..1 floats; PPTX wants 0..255 ints."""
                if c is None:
                    return None
                try:
                    r, g, b = (max(0, min(255, int(v * 255))) for v in c[:3])
                    return RGBColor(r, g, b)
                except (TypeError, ValueError):
                    return None

            try:
                prs = Presentation()
                prs.slide_width = Emu(int(slide_w_pt * 12700))
                prs.slide_height = Emu(int(slide_h_pt * 12700))
                blank = prs.slide_layouts[6]
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    slide = prs.slides.add_slide(blank)

                    # ── Phase 1: vector drawings (filled rects, lines).
                    # Headers, banner bars, card backgrounds, separators
                    # in slide-builder PDFs are vector drawings, not text
                    # or images. Without this phase the slide looked
                    # "naked" — only text floating on a white background.
                    # Added FIRST so subsequent text/image shapes land on
                    # top in PowerPoint's z-order.
                    try:
                        drawings = page.get_drawings()
                    except Exception:
                        drawings = []
                    for d in drawings:
                        items = d.get("items") or []
                        rect = d.get("rect")
                        fill = _rgb(d.get("fill"))
                        stroke = _rgb(d.get("stroke"))
                        if rect is None:
                            continue
                        x0, y0, x1, y1 = rect
                        x0 = max(0, x0); y0 = max(0, y0)
                        x1 = min(slide_w_pt, x1); y1 = min(slide_h_pt, y1)
                        w = x1 - x0; h = y1 - y0
                        if w <= 0 or h <= 0:
                            continue
                        # Sub-point noise (anti-aliasing or clipping
                        # artifacts) — skip
                        if w * h < 0.25:
                            continue
                        kinds = {it[0] for it in items}
                        # Render *any* drawing with a fill colour as a
                        # rectangle at its bbox. The earlier strict
                        # filter (only `re`+`l` paths) silently
                        # rejected rounded-rect cards drawn with `c`
                        # curves at the corners — common in modern
                        # slide-builder PDFs (Genially, Canva, etc.).
                        # Map drawings whose path contains Bezier
                        # curves to ROUNDED_RECTANGLE for better
                        # fidelity; sharp paths (only `re`/`l`) stay
                        # RECTANGLE.
                        # Stroked-only thin shapes still become a thin
                        # filled rect in the stroke colour.
                        is_filled = fill is not None and items
                        is_line = (fill is None and stroke is not None
                                   and items
                                   and w * h <= 4 * max(w, h))
                        if not (is_filled or is_line):
                            continue
                        is_rounded = is_filled and bool({"c", "qu"} & kinds)
                        shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE
                                      if is_rounded else MSO_SHAPE.RECTANGLE)
                        try:
                            shape = slide.shapes.add_shape(
                                shape_type,
                                Emu(int(x0 * 12700)),
                                Emu(int(y0 * 12700)),
                                Emu(int(w * 12700)),
                                Emu(int(h * 12700)))
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill if is_filled else stroke
                            try:
                                # Border off — PDF drawings are usually
                                # fill-only; the stroke (if any) is the
                                # same colour or absent.
                                shape.line.fill.background()
                            except Exception:
                                pass  # noqa: S110
                            # Strip the auto-generated <p:style> block
                            # python-pptx adds to every add_shape call.
                            # That block carries `effectRef idx="2"` and
                            # `fillRef idx="3"` pointing to theme presets;
                            # PowerPoint silently applies those *on top of*
                            # the explicit `<a:solidFill>` we just set on
                            # some slides, leading to invisible / wrongly
                            # tinted shapes (observed: slide 6 of the user's
                            # UFCD deck rendered as blank white in
                            # PowerPoint despite all 13 rects being in the
                            # XML with valid coords and colours). Removing
                            # the style block forces PowerPoint to use only
                            # `<p:spPr>`, which is what we control.
                            try:
                                style_el = shape._element.find(qn("p:style"))
                                if style_el is not None:
                                    shape._element.remove(style_el)
                            except Exception:
                                pass  # noqa: S110
                        except Exception:
                            # Some MSO_SHAPE / fill combinations fail
                            # silently in older python-pptx; skip the
                            # drawing rather than abort the whole slide.
                            pass  # noqa: S110

                    # ── Phase 2: text / image blocks.
                    # Editable extraction: image blocks become picture
                    # shapes; each *line* of a text block becomes its
                    # own textbox at the line's bbox so PowerPoint
                    # doesn't re-flow lines (one textbox per block was
                    # wrong — multi-line blocks lost their original
                    # vertical positions when PowerPoint re-laid out
                    # the text). Spans within the same line stay as
                    # runs in a single paragraph so inline formatting
                    # (bold/italic mid-sentence) is preserved.
                    blocks = page.get_text("dict").get("blocks", [])
                    for block in blocks:
                        if block.get("type") == 1:  # image block
                            bbox = block.get("bbox")
                            img_data = block.get("image")
                            if not (bbox and img_data):
                                continue
                            ix0, iy0, ix1, iy1 = bbox
                            ix0 = max(0, ix0); iy0 = max(0, iy0)
                            ix1 = min(slide_w_pt, ix1); iy1 = min(slide_h_pt, iy1)
                            iw = ix1 - ix0; ih = iy1 - iy0
                            if iw <= 0 or ih <= 0:
                                continue
                            try:
                                slide.shapes.add_picture(
                                    io.BytesIO(img_data),
                                    Emu(int(ix0 * 12700)),
                                    Emu(int(iy0 * 12700)),
                                    Emu(int(iw * 12700)),
                                    Emu(int(ih * 12700)))
                            except Exception:
                                # Image format not supported by python-pptx
                                # (rare formats like JBIG2). Skip the
                                # block — the rest of the slide is still
                                # produced.
                                pass
                            continue

                        for line in block.get("lines", []):
                            spans = [s for s in line.get("spans", [])
                                     if s.get("text", "")]
                            if not spans:
                                continue
                            lbb = line.get("bbox")
                            if not lbb:
                                continue
                            lx0, ly0, lx1, ly1 = lbb
                            # Pad horizontally so PowerPoint's slightly
                            # different glyph metrics don't clip the
                            # last character; vertically use the line
                            # bbox as-is to keep baselines aligned.
                            lx0 = max(0, lx0 - 1)
                            ly0 = max(0, ly0)
                            lx1 = min(slide_w_pt, lx1 + 4)
                            ly1 = min(slide_h_pt, ly1)
                            lw = lx1 - lx0; lh = ly1 - ly0
                            if lw <= 0 or lh <= 0:
                                continue
                            try:
                                tb = slide.shapes.add_textbox(
                                    Emu(int(lx0 * 12700)),
                                    Emu(int(ly0 * 12700)),
                                    Emu(int(lw * 12700)),
                                    Emu(int(lh * 12700)))
                            except Exception:
                                continue
                            tf = tb.text_frame
                            tf.word_wrap = False  # one line — no wrap
                            # Zero internal padding so the textbox bbox
                            # matches the PDF line bbox more faithfully.
                            # Older python-pptx versions reject Emu(0)
                            # for margins; tolerate that.
                            for attr in ("margin_left", "margin_right",
                                         "margin_top", "margin_bottom"):
                                try: setattr(tf, attr, 0)
                                except Exception: pass  # noqa: S110

                            para = tf.paragraphs[0]
                            first_run = True
                            for span in spans:
                                text = span["text"]
                                # Reuse the auto-created empty run for
                                # the first span; add new runs after.
                                if first_run and len(para.runs) > 0:
                                    run = para.runs[0]
                                else:
                                    run = para.add_run()
                                first_run = False
                                run.text = text
                                size = span.get("size", 12)
                                # Pt() rejects negative or non-numeric
                                # sizes; if the PDF span had a junk
                                # value, fall back to the theme default.
                                try: run.font.size = Pt(size)
                                except Exception: pass  # noqa: S110
                                flags = span.get("flags", 0)
                                run.font.bold = bool(flags & 16)
                                run.font.italic = bool(flags & 2)
                                # PDF subset fonts come prefixed with
                                # 6 random caps + "+", strip them so
                                # PowerPoint can substitute the
                                # canonical face: e.g. "ABCDEF+Arial"
                                # → "Arial".
                                font_name = span.get("font", "") or ""
                                if len(font_name) > 7 and font_name[6] == "+":
                                    font_name = font_name[7:]
                                if font_name:
                                    try: run.font.name = font_name
                                    except Exception: pass  # noqa: S110
                                color = span.get("color", 0)
                                if color:
                                    try:
                                        run.font.color.rgb = RGBColor(
                                            (color >> 16) & 0xFF,
                                            (color >> 8) & 0xFF,
                                            color & 0xFF)
                                    except Exception:
                                        # Some run types reject explicit
                                        # color (e.g. inside placeholder
                                        # layouts); leave the default.
                                        pass
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                prs.save(out_path)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(t("tool.convert.status.pptx_done", path=out_path))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_pptx", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → XLSX ──────────────────────────────────────────────────────

    def _convert_xlsx(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.xlsx")
        if not out_path:
            return
        try:
            from openpyxl import Workbook  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_xlsx"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            show_error(self, e)
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            from openpyxl import Workbook
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                wb = Workbook()
                wb.remove(wb.active)
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    ws = wb.create_sheet(
                        title=t("tool.convert.xlsx.sheet_name", n=i + 1))
                    blocks = page.get_text("blocks")
                    for row_idx, block in enumerate(blocks):
                        if block[6] != 0:  # skip image blocks
                            continue
                        text = _clean(block[4].strip())
                        if text:
                            cells = [c.strip() for c in text.replace("\t", "|").split("|") if c.strip()]
                            if not cells:
                                cells = [text]
                            for col_idx, cell in enumerate(cells):
                                ws.cell(row=row_idx + 1, column=col_idx + 1, value=cell)
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                wb.save(out_path)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(t("tool.convert.status.xlsx_done", path=out_path))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_xlsx", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → HTML ──────────────────────────────────────────────────────

    def _convert_html(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.html")
        if not out_path:
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            show_error(self, e)
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                parts = [
                    "<!DOCTYPE html>",
                    '<html lang="en"><head><meta charset="UTF-8">',
                    f"<title>{os.path.basename(pdf_path)}</title>",
                    "<style>body{font-family:sans-serif;max-width:800px;margin:0 auto;padding:20px;}"
                    ".page{margin-bottom:40px;padding-bottom:20px;border-bottom:1px solid #ccc;}</style>",
                    "</head><body>",
                ]
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    parts.append('<div class="page">')
                    blocks = page.get_text("dict")["blocks"]
                    for block in blocks:
                        if block.get("type") != 0:
                            continue
                        for line in block.get("lines", []):
                            spans_html = ""
                            for span in line.get("spans", []):
                                text = span["text"]
                                if not text.strip():
                                    continue
                                flags = span.get("flags", 0)
                                bold = flags & 16
                                italic = flags & 2
                                tag_text = _clean(text).replace("&", "&amp;").replace("<", "&lt;")
                                if bold:
                                    tag_text = f"<strong>{tag_text}</strong>"
                                if italic:
                                    tag_text = f"<em>{tag_text}</em>"
                                spans_html += tag_text
                            if spans_html:
                                avg_size = max(s.get("size", 12) for s in line.get("spans", [{"size": 12}]))
                                if avg_size >= 18:
                                    parts.append(f"<h1>{spans_html}</h1>")
                                elif avg_size >= 15:
                                    parts.append(f"<h2>{spans_html}</h2>")
                                elif avg_size >= 13:
                                    parts.append(f"<h3>{spans_html}</h3>")
                                else:
                                    parts.append(f"<p>{spans_html}</p>")
                    parts.append("</div>")
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                parts.append("</body></html>")
                if worker.is_cancelled():
                    return None
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(parts))
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(t("tool.convert.status.html_done", path=out_path))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_html", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)

    # ── PDF → EPUB ──────────────────────────────────────────────────────

    def _convert_epub(self, pdf_path: str):
        out_path = self._resolve_output_file(self._drop_file, pdf_path,
                                             filter_key="file_filter.epub")
        if not out_path:
            return
        try:
            from ebooklib import epub  # noqa: F401
        except ImportError:
            QMessageBox.critical(self, t("msg.missing_dep"), t("tool.convert.dep_epub"))
            return
        try:
            with self._open_fitz(pdf_path) as _probe:
                total = _probe.page_count
        except Exception as e:
            show_error(self, e)
            return
        if total == 0:
            QMessageBox.warning(self, t("msg.warning"), t("msg.select_valid_pdf"))
            return
        pwd = self._pdf_password

        def do_work(worker):
            import fitz
            from ebooklib import epub
            doc = fitz.open(pdf_path)
            if doc.needs_pass and pwd:
                doc.authenticate(pwd)
            try:
                book = epub.EpubBook()
                book.set_identifier(f"pdfapps-{os.path.basename(pdf_path)}")
                book.set_title(os.path.splitext(os.path.basename(pdf_path))[0])
                book.set_language("en")
                chapters = []
                for i, page in enumerate(doc):
                    if worker.is_cancelled():
                        return None
                    page_title = t("tool.convert.epub.page_title", n=i + 1)
                    ch = epub.EpubHtml(title=page_title,
                                       file_name=f"page_{i+1}.xhtml")
                    text = page.get_text()
                    paragraphs = [f"<p>{_clean(p)}</p>" for p in text.split("\n") if p.strip()]
                    ch.content = f"<html><body><h2>{page_title}</h2>{''.join(paragraphs)}</body></html>"
                    book.add_item(ch)
                    chapters.append(ch)
                    worker.progress.emit(i, f"{i + 1}/{total}…")
                if worker.is_cancelled():
                    return None
                book.add_item(epub.EpubNcx())
                book.add_item(epub.EpubNav())
                book.spine = ["nav"] + chapters
                book.toc = chapters
                epub.write_epub(out_path, book)
            finally:
                doc.close()
            return total

        def on_done(result):
            self.lbl_result.setText(f"  → {os.path.basename(out_path)}")
            self._status(t("tool.convert.status.epub_done", path=out_path))
            QMessageBox.information(self, t("msg.done"),
                                    t("tool.convert.done_epub", path=out_path))

        self._run_background(do_work, total, t("tool.convert.converting"),
                             on_done=on_done)
